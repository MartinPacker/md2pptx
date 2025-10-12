"""
runPython
"""

version = "0.11"

import csv
from pptx.chart.data import CategoryChartData
from pptx.oxml.xmlchemy import OxmlElement, serialize_for_reading
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from colour import setColour, parseColour
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from paragraph import *
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE

from media import *

import globals

class RunPython:
    def __init__(
        self,
    ):
        pass

    # Execute the lines of code passed in
    def run(self, prs, slide, renderingRectangle, codeLines, codeType):
        concatenatedCodeLines = "\n".join(codeLines)
        exec(concatenatedCodeLines)

    def runFromFile(self, filename, prs, slide, renderingRectangle):
        exec(open(filename).read())
    

    # Helper function for run-python
    def readCSV(filename):
        my_csv = []
        with open(filename, 'r') as csvfile:
            chart_reader = csv.reader(csvfile, quoting = csv.QUOTE_NONNUMERIC)
            for row in chart_reader:
                my_csv.append(row)

        return my_csv
    
    def filterRows(my_array, filterFunction):
        my_array2 = []
        for rowNumber, row in enumerate(my_array):
            if filterFunction(rowNumber, row):
                my_array2.append(row)
        
        return my_array2
        

    def transposeArray(chart_array):
        return list(map(list, zip(*chart_array)))
    
    def makeChartData(chart_array, seriesIsColumn = True, columns = None):

        chart_data = CategoryChartData()
        
        if columns is not None:
            chart_array2 = []
            for rowNumber, row in enumerate(chart_array):
                chart_row = []
                for column in columns:
                    chart_row.append(row[column])
                
                chart_array2.append(chart_row)
            chart_array = chart_array2
            print(pptx.enum.chart)


        if seriesIsColumn:
            # Transpose input data
            chart_array = RunPython.transposeArray(chart_array)
        
        # x values
        chart_data.categories = chart_array[0][1:]

        # Series
        for rowNumber, row in enumerate(chart_array[1:]):
            chart_data.add_series(row[0],row[1:])

        return chart_data
  
    # Helper function to make a chart. The result can be further manipulated
    def makeChart(slide,
        chart_type,
        renderingRectangle,
        chart_data,
        title = None,
        legendPosition = None):    
        c = slide.shapes.add_chart(
            chart_type,
            renderingRectangle.left,
            renderingRectangle.top,
            renderingRectangle.width, 
            renderingRectangle.height,
            chart_data
        )
          
        chart = c.chart
        
        if title is not None:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
          
        if legendPosition is not None:
            chart.has_legend = True
            chart.legend.position = legendPosition
            chart.legend.include_in_layout = False


        return c

    # Helper routine to make a table. The result can be further manipulated
    def makeTable(slide,
        renderingRectangle,
        table_array):

        height = len(table_array)
        width = len(table_array[0])

        t = slide.shapes.add_table(height,width,
            renderingRectangle.left,
            renderingRectangle.top,
            renderingRectangle.width,
            renderingRectangle.height)

        table = t.table
        
        for i in range(height):
            for j in range(width):
               c = table.cell(i, j)
               c.text = str(table_array[i][j])

        return t

    def applyCellFillRGB(table, row, column, red, green, blue):
        cff = table.table.cell(row, column).fill
        cff.solid()
        cff.fore_color.rgb = RGBColor(red, green, blue)

    def applyCellListFillRGB(table, cellList, red, green, blue):
        for row, column in cellList:
            RunPython.applyCellFillRGB(table, row, column, red, green, blue)
            
    def alignTableCellText(tableFrame, rowNumber, columnNumber, alignment, paragraphNumber = None):
        # Get the cell's text_frame
        tableCellFrame = tableFrame.table.cell(rowNumber, columnNumber).text_frame

        if paragraphNumber == None:
            # Iterate over the cell's paaragraph's, aligning right
            for p in tableCellFrame.paragraphs:
                p.alignment = alignment
        else:
            tableCellFrame.paragraphs[paragraphNumber].alignment = alignment

    def makeDrawnShape(slide, vertices, fill = False, text = None, textColor = None, fillColor = None, closed = True):
        ffBuilder = slide.shapes.build_freeform(*vertices[0], True)
    
        ffBuilder.add_line_segments(vertices[1:], close = closed)
    
        s = ffBuilder.convert_to_shape()
        
        if text is not None:
            s.text = text
            p = s.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            if textColor is None:
                setColour(p.font.color, parseColour('#000000'))
            else:
                setColour(p.font.color, parseColour(textColor))
    
        if fill:
            s.fill.solid()
            if fillColor is not None:
                setColour(s.fill.fore_color, parseColour(fillColor))
        
        return s

    def doChecklistChecks(placeholder, checklist, colourChecks = False):
        slide = placeholder._parent._parent
        tf = placeholder.text_frame
        paras = tf.paragraphs
        for paraNumber, para in enumerate(paras):
            # Save original font size
            originalFontSize = para.font.size
            
            # Save original indentation level
            level = para.level
            
            # Remove the original pPr element
            para._element.remove(para._element.getchildren()[0])
    
            xml = ''
            
            # Note the level insertion
            xml += f'<a:pPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" marL="{Inches(0.25) * level}" indent="{Inches(0.33)}" lvl="{level}">'


            if checklist[paraNumber] == "Unset":
                # Set the bullet to an empty square
                image_part, rId = createMediaRel(slide, "unset-black.png")
                
                # Following is non-graphic bullet
                # xml += '<a:buFont xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="Wingdings" pitchFamily="2" charset="2"/>'
                # xml += '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="o"/>'

            elif checklist[paraNumber] == "Yes":
                # Set the bullet to a square with a tick
                if colourChecks:
                    image_part, rId = createMediaRel(slide, "tick-colour.png")
                else:
                    image_part, rId = createMediaRel(slide, "tick-black.png")
                
                # Following is non-graphic colouring
                # xml += '<a:buClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                # xml += '<a:srgbClr val="00FF00" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />'
                # xml += '</a:buClr>'

                # Following is non-graphic bullet
                # xml += '<a:buFont xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="Wingdings 2" pitchFamily="2" charset="2"/>'
                # xml += '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="R"/>'

            elif checklist[paraNumber] == "Maybe":
                # Bullet set to maybe
                if colourChecks:
                    image_part, rId = createMediaRel(slide, "query-colour.png")
                else:
                    image_part, rId = createMediaRel(slide, "query-black.png")
                
                # Following is non-graphic colouring
                # xml += '<a:buClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                # xml += '<a:srgbClr val="FFA500" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />'
                # xml += '</a:buClr>'

                # Following is non-graphic bullet
                # xml += '<a:buFont xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="Wingdings" pitchFamily="2" charset="2"/>'
                # xml += '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="⍰"/>'

            elif checklist[paraNumber] == "Partial":
                # Bullet set to partial
                if colourChecks:
                    image_part, rId = createMediaRel(slide, "partial-colour.png")
                else:
                    image_part, rId = createMediaRel(slide, "partial-black.png")
                
                # Following is non-graphic colouring
                # xml += '<a:buClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                # xml += '<a:srgbClr val="0000FF" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />'
                # xml += '</a:buClr>'

                # Following is non-graphic bullet
                # xml += '<a:buFont xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="Courier" pitchFamily="2" charset="2"/>'
                # xml += '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="▃"/>'

            else:
                # Set the bullet to a square with a cross
                if colourChecks:
                    image_part, rId = createMediaRel(slide, "cross-colour.png")
                else:
                    image_part, rId = createMediaRel(slide, "cross-black.png")
                
                
                # Following is non-graphic colouring
                # xml += '<a:buClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                # xml += '<a:srgbClr val="FF0000" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />'
                # xml += '</a:buClr>'

                # Following is non-graphic bullet
                # xml += '<a:buFont xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="Wingdings 2" pitchFamily="2" charset="2"/>'
                # xml += '<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="T"/>'
     
            xml += '<a:buBlip xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            xml += f'    <a:blip xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" r:embed="{rId}"/>'
            xml += '</a:buBlip>'


            xml += '</a:pPr>'
    
            # Parse this XML
            parsed_xml = parse_xml(xml)
    
            # Insert the parsed XML fragment as a child of the pPr element
            para._element.insert(0, parsed_xml)
    
            # Restore original font size
            para.font.size = originalFontSize

        return placeholder

    def makeChecklist(placeholder, checklist, checkTextIndex = 0, checkMarkIndex = 1, levelIndex = 2, colourChecks = False):
        checkMarks = []
        
        for paraNumber, checklistItem in enumerate(checklist):
            checkMarks.append(checklistItem[checkMarkIndex])

            if paraNumber == 0:
               para = placeholder.text_frame.paragraphs[0]
            else:
               para = placeholder.text_frame.add_paragraph()

            addFormattedText(para, checklistItem[checkTextIndex])
            
            # The actual level setting in the surviving XML is done by doChecklistChecks
            # The below sets the level as a hint for doChecklistChecks to work with
            if len(checklistItem) > levelIndex:
                try:
                    level = int(checklistItem[levelIndex])
                    
                    para.level = level - 1
                except ValueError:
                    para.level = 0
            else: para.level = 0

        RunPython.doChecklistChecks(placeholder, checkMarks, colourChecks)
        
        return placeholder

    def testForValues(cellValue, testValues):
        if isinstance(testValues, list):
            # Work through list of possible values
            for testValue in testValues:
                if cellValue == testValue:
                    return True

            return False
        else:
            # Single string value to check for
            return cellValue == testValues

    def makeTruthy(table_array, columnNumber = 1, trueValues = "yes", falseValues = "no", unsetValues = "", maybeValues = "maybe", partialValues = "partial"):
        for row in table_array:
            # Test is case insensitive
            cellValue = row[columnNumber].lower()

            if RunPython.testForValues(cellValue, unsetValues):
                row[columnNumber] = "Unset"

            elif RunPython.testForValues(cellValue, trueValues):
                row[columnNumber] = "Yes"

            elif RunPython.testForValues(cellValue, maybeValues):
                row[columnNumber] = "Maybe"

            elif RunPython.testForValues(cellValue, partialValues):
                row[columnNumber] = "Partial"

            elif RunPython.testForValues(cellValue, falseValues):
                row[columnNumber] = "No"

        return table_array

    def ensureTextbox(slide, renderingRectangle, shapeIndex = None):
        if shapeIndex is not None:
            if len(slide.shapes) < shapeIndex + 1:
                # Need to create a text box as shape index too high
                newShape = slide.shapes.add_textbox(
                    renderingRectangle.left,
                    renderingRectangle.top,
                    renderingRectangle.width,
                    renderingRectangle.height
                )

                # Return new shape
                return newShape
            else:
                # shape Index is valid
                if (theShape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX) | (theShape.placeholder_format.type == PP_PLACEHOLDER.OBJECT):
                    # shape is a text box so return it
                    return slide.shapes[shapeIndex]
                else:
                    # Shape isn't a text box so create one
                    newShape = slide.shapes.add_textbox(renderingRectangle.left, renderingRectangle.top, renderingRectangle.width, renderingRectangle.height)

                    # Return new shape
                    return newShape

        # shapeIndex wasn't set so potentially search for the last text box
        if len(slide.shapes) < 2:
            # Need to create a text box as only shape is presumed to be a title
            newShape = slide.shapes.add_textbox(renderingRectangle.left, renderingRectangle.top, renderingRectangle.width, renderingRectangle.height)

            # Return new shape
            return newShape
        
        # Search for last text box
        for shapeIndex, theShape in reversed(list(enumerate(slide.shapes))):
            if (theShape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX) | (theShape.placeholder_format.type == PP_PLACEHOLDER.OBJECT):
                return theShape


        # Need to create a text box none found - other than perhaps at Index 0 (title)
        newShape = slide.shapes.add_textbox(renderingRectangle.left, renderingRectangle.top, renderingRectangle.width, renderingRectangle.height)

        # Return new shape's index as presumed to be the text box we want
        return newShape
        
    def checklistFromCSV(slide, renderingRectangle, filename, shapeIndex = None, colourChecks = False):
        # Read in CSV and turn second column into "truthy" values
        myChecklist = RunPython.makeTruthy(RunPython.readCSV(filename),
            1,
            trueValues = ["yes", "x", "y"],
            falseValues = ["no", "n"],
            unsetValues = ["", " "],
            maybeValues = ["maybe", "m", "?"],
            partialValues = ["partial", "p"],
        )

        # Ensure we have a placeholder - whether first or second or specified
        textShape = RunPython.ensureTextbox(slide, renderingRectangle, shapeIndex)

        # Make the checklist in this placeholder from the imported file
        RunPython.makeChecklist(textShape, myChecklist, 0, 1, 2, colourChecks)
        
        return textShape

    def removeBullet(theShape, paragraphNumber):
        removeBullet(theShape.text_frame.paragraphs[paragraphNumber])
        
        return theShape
    
    def removeBullets(theShape):
        for p in theShape.text_frame.paragraphs:
            removeBullet(p)

        return theShape

    def removeSelectedBullets(theShape, removalArray):
        removeSelectedBullets(theShape.text_frame, removalArray)

        return theShape

    def getParagraphs(slide, wantedParagraphs = []):
        return getParagraphs(slide, wantedParagraphs)
        
    def doAnnotations(slide, annotationList, lineWidth = None, shapeWidth = None):
        for annotation in annotationList:
            x = Inches(float(annotation[0]))
            y = Inches(float(annotation[1]))
            w = Inches(float(annotation[2]))
            h = Inches(float(annotation[3]))
            
            text = annotation[4]
            
            if text in [
                "-",
                "<-",
                "->",
                "<->",
                "=",
                "<=",
                "=>",
                "<=>",
            ]:
                # Draw an line from x, y to x+w, y+h
                c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y + h)
                
                if text != "-":
                    # Will need an a:ln element
                    
                    # Find the spPr element to hang this off
                    for element in c._element.getchildren():
                        if element.tag == "{http://schemas.openxmlformats.org/presentationml/2006/main}spPr":
                            spPr = element
                            break
                    if "=" in text:
                        cmpd = "dbl"
                    else:
                        cmpd = "sng"

                    xml = '<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" cmpd="' + cmpd + '">'
                        
                    if "<" in text:
                        xml += '  <a:headEnd type="triangle" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />'
                        
                    if ">" in text:
                        xml += '  <a:tailEnd type="triangle" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" />'
                        
                    xml += '</a:ln>'
                        
                    # Parse this XML
                    parsed_xml = parse_xml(xml)
    
                    # Insert the parsed XML fragment as a child of the pPr element
                    spPr.append(parsed_xml)
                if lineWidth is not None:
                    c.line.width = Pt(float(lineWidth))
                        
                if len(annotation) > 5:
                    toColour = c.line.color
                    setColour(toColour, parseColour(annotation[5]))

            elif text[0] == "!":
                filename = text[1:]
                slide.shapes.add_picture(filename, x, y, w, h)

            elif text in [
                "[]",
                "()",
                "[-]",
                "(-)",
                "[=]",
                "(=)",
                "o",
                "O",
            ]:
                if text in [
                    "[]",
                    "[-]",
                    "[=]",
                ]:
                    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
                    
                elif text in [
                    "o",
                    "O",
                ]:
                    b = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
                    
                else:
                    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)

                if len(annotation) > 5:
                    b.text = annotation[5]
                    f = b.text_frame
                    p = f.paragraphs[0].alignment = PP_ALIGN.CENTER

                if len(annotation) > 6:
                    # Foreground colour
                    foreColour = annotation[6]
                    if foreColour != "":
                        toColour = b.text_frame.paragraphs[0].runs[0].font.color
                        setColour(toColour, parseColour(foreColour))
                    
                    if len(annotation) > 7:
                        # Background colour
                        backColour = annotation[7]
                        if backColour != "":
                            b.fill.solid()
                            toColour = b.fill.fore_color
                            setColour(toColour, parseColour(annotation[7]))
                    
                if shapeWidth is not None:
                    b.line.width = Pt(float(shapeWidth))
                
                if ("=" in text) | (text == "O"):
                    be = b._element
                    ln= b.get_or_add_ln()
                    ln.set("cmpd", "dbl")
            else:
                t = slide.shapes.add_textbox(x, y, w, h)
                t.text = text
                if len(annotation) > 5:
                    toColour = t.text_frame.paragraphs[0].runs[0].font.color
                    setColour(toColour, parseColour(annotation[5]))

    def annotationsFromCSV(slide, filename, lineWidth = None, shapeWidth = None):
        annotations = RunPython.readCSV(filename)
        
        RunPython.doAnnotations(slide, annotations, lineWidth, shapeWidth)