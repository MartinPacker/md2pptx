"""
funnel
"""

myVersion = "0.0"

__version__ = myVersion

import csv
import io
from rectangle import Rectangle
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from colour import setColour
from symbols import resolveSymbols

def massageFunnelText(text):
    fragment = ""
    for c in text:
        if ord(c) == 236:
            fragment = fragment + "<"

        elif ord(c) == 237:
            fragment = fragment + ">"

        else:
            fragment = fragment + c
            
    return fragment

class Funnel:
    def __init__(
        self,
    ):
        pass
    
    def makeFunnel(
        self, 
        slide, 
        renderingRectangle, 
        funnelParts,
        partColours,
        codeType,
        funnelBorderColour,
        funnelTitleColour,
        funnelTextColour,
        funnelLabelsPercent,
    ):
        funnelLabelsProportion = funnelLabelsPercent / 100
        
        # Define labels rectangle
        funnelLabelsRectangle = Rectangle(
            renderingRectangle.top,
            renderingRectangle.left,
            int(renderingRectangle.height * funnelLabelsProportion),
            renderingRectangle.width,
        )

        # Define funnel body rectangle
        funnelBodyRectangle = Rectangle(
            renderingRectangle.top + int(renderingRectangle.height * funnelLabelsProportion),
            renderingRectangle.left,
            int(renderingRectangle.height * (1 - funnelLabelsProportion)),
            renderingRectangle.width,
        )

        tipHeight = funnelBodyRectangle.height / 3
        
        
        partColourCount = len(partColours)
    
        funnelPartRows = [r for r in csv.reader(io.StringIO(str.join("\n", funnelParts)), escapechar = "\\", skipinitialspace = True)]

        funnelPartCount = len(funnelPartRows)
        
        # Build lists of labels and Body
        funnelLabels = []
        funnelBody = []
        
        for row in funnelPartRows:
            cell1 = row[0].strip()
            if len(row) == 0:
                funnelLabels.append("")
                funnelBody.append("")
            else:
                funnelLabels.append(cell1)
                if len(row) == 1:
                    funnelBody.append("")
                else:
                    cell2 = row[1].strip()
                    funnelBody.append(cell2)
                
        partWidth = renderingRectangle.width / funnelPartCount

        # Create the labels
        for l, label in enumerate(funnelLabels):
            tb = slide.shapes.add_textbox(
                funnelLabelsRectangle.left + l * partWidth,
                funnelLabelsRectangle.top,
                partWidth,
                funnelLabelsRectangle.height,
            )
            
            tb.text = massageFunnelText(resolveSymbols(label.replace("<br/>","\n")))
            for p in tb.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                if funnelTitleColour != ("None", ""):
                    setColour(p.font.color, funnelTitleColour)

        # Create the parts of the funnel§
        for b, body in enumerate(funnelBody):
            partLeft = funnelBodyRectangle.left + b * partWidth
            partRight = partLeft + partWidth
            
            if b == funnelPartCount - 1:
                leftSpaceAboveBelow = (funnelBodyRectangle.height - tipHeight) / 2
                partTopLeft = funnelBodyRectangle.top + leftSpaceAboveBelow
                partTopRight = partTopLeft
                partBottomLeft = funnelBodyRectangle.top + funnelBodyRectangle.height - leftSpaceAboveBelow
                partBottomRight = partBottomLeft
            else:
                leftSpaceAboveBelow = (funnelBodyRectangle.height - tipHeight) / 2 * b / (funnelPartCount - 1)
                rightSpaceAboveBelow = (funnelBodyRectangle.height - tipHeight) / 2 * (b + 1) / (funnelPartCount - 1)
                partTopLeft = funnelBodyRectangle.top + leftSpaceAboveBelow
                partTopRight = funnelBodyRectangle.top + rightSpaceAboveBelow
                partBottomLeft = funnelBodyRectangle.top + funnelBodyRectangle.height - leftSpaceAboveBelow
                partBottomRight = funnelBodyRectangle.top + funnelBodyRectangle.height - rightSpaceAboveBelow
            
            # Start shape builder with first point
            ffBuilder = slide.shapes.build_freeform(
                partLeft, 
                partTopLeft
            )
            
            ffBuilder.add_line_segments(
                [
                    (partLeft , partBottomLeft), 
                    (partRight, partBottomRight),
                    (partRight, partTopRight),
                ],
                close = True
            )
    
            s = ffBuilder.convert_to_shape()
            s.text = massageFunnelText(resolveSymbols(body.replace("<br/>","\n")))
            
            for p in s.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                if funnelTextColour != ("None", ""):
                    setColour(p.font.color, funnelTextColour)

            s.fill.solid()
            
            partColourType, partColourValue = partColours[b % partColourCount]
            if partColourType == "Theme":
                s.fill.fore_color.theme_color = partColourValue
            else:
                s.fill.fore_color.rgb = RGBColor.from_string(partColourValue[1:])
            
            if funnelBorderColour != ("None", ""):
                setColour(s.line.color, funnelBorderColour)



