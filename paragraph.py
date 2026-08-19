"""
paragraph
"""

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
from lxml import etree
from pptx.oxml.xmlchemy import OxmlElement
import re
import sys
from pptx.dml.color import RGBColor, MSO_THEME_COLOR
from pptx.util import Pt

import globals
import pptx_math
from processingOptions import *
from symbols import resolveSymbols
from colour import *

# Inline maths: $`LaTeX`$ inside body text. Block maths already had a handler
# (``` math fenced blocks -> pptx_math.PptxMath), but there was no way to put a
# symbol inside a sentence, so authors had to spell "L_patch" in plain text.
#
# The delimiters are GitHub's, and deliberately not a bare "$". parseText below
# rewrites the line before tokenising it, and the tokeniser then drops "*" and
# "[" whenever they match none of its cases - even inside a code span - so LaTeX
# left in the line does not survive. Requiring a backtick against the dollar
# also keeps "costs $5 to $10" out of it.
#
# The LaTeX may not contain a backtick. GitHub's inner delimiters are a code
# span, which cannot contain one either, and allowing it lets an unclosed opener
# pair with the closing delimiter of the *next* formula and swallow the prose in
# between.
inlineMathRegex = re.compile(r"(?<!\\)\$`([^`]+)`\$")

# Each formula is lifted out of the line and replaced by this one character.
# FDD0-FDEF are permanent Unicode noncharacters, so they cannot occur in the
# author's text. FDD0-FDE3 are already taken by the constructs parseText handles.
mathSentinel = u"\uFDE4"

# Anything in this range is one of the sentinels above, already substituted into
# the text before parseText saw it. The section title, presentation title and
# presentation subtitle routes resolveSymbols() their text before handing it over,
# so a "\`" or an entity reference inside a formula arrives as a noncharacter.
# Converting that buries an unreadable character in the OMML, so such a formula
# is left as its source.
substitutedSentinelRegex = re.compile("[\uFDD0-\uFDEF]")

# The tokeniser states in which a formula can become a fragment of its own.
# Everything else builds its fragment up as structure that a fragment boundary
# would break - link text, a footnote reference - so those keep the literal
# source instead. Listing what is allowed rather than what is not means a state
# added later degrades to visible text rather than corrupting the fragment.
mathFragmentStates = [
    "N", "I", "B1", "B2", "C",
    "CMRep", "CMHig", "CMCom", "CMDel", "CMAdd",
    "Ins", "Del", "Sub", "Sup",
]

# Compiling the XSLT is slow, so keep one inserter per mathxsl option value.
_mathInserterCache = {}


def getMathInserter():
    """MathInserter for the resolved stylesheet, or None if there is not one.

    _resolve_mathxsl_path takes the mathxsl option when it names a file and
    falls back to the copy beside pptx_math.py, so inline maths becomes
    available on exactly the same terms as a ``` math block.
    """
    configured = globals.processingOptions.getCurrentOption("mathxsl")
    if configured not in _mathInserterCache:
        try:
            _mathInserterCache[configured] = pptx_math.MathInserter(
                pptx_math._resolve_mathxsl_path(configured)
            )
        except Exception as e:
            sys.stderr.write(f"Inline math skipped: {e}\n")
            _mathInserterCache[configured] = None
    return _mathInserterCache[configured]


def inlineMathSource(latex):
    """The literal "$`...`$" the author wrote, rebuilt from its LaTeX."""
    return "$`" + latex + "`$"


def emitInlineMath(p, latex):
    """Append `latex` to paragraph `p` as one inline OMML equation.

    pptx_math.inject_inline_math appends to the paragraph, so calling it
    between add_run() calls keeps the reading order.

    A formula that will not convert is left as its literal "$`...`$" source
    rather than dropped: a symbol silently missing from the middle of a
    sentence is far harder to notice than a stray dollar sign.
    """
    inserter = getMathInserter()
    omath = None
    if inserter is not None:
        try:
            omath = inserter.make_inline_omml(latex)
        except Exception as e:
            sys.stderr.write(f"Inline math conversion failed for '{latex}': {e}\n")

    if omath is None:
        p.add_run().text = inlineMathSource(latex)
    else:
        pptx_math.inject_inline_math(p._p, omath)


# Following functions are workarounds for python-pptx not having these functions for the font object
def setSubscript(font):
    if font.size is None:
        font._element.set("baseline", "-50000")
        return
    
    if font.size < Pt(24):
        font._element.set("baseline", "-50000")
    else:
        font._element.set("baseline", "-25000")


def setSuperscript(font):
    if font.size is None:
        font._element.set("baseline", "60000")
        return
    
    if font.size < Pt(24):
        font._element.set("baseline", "60000")
    else:
        font._element.set("baseline", "30000")


def setStrikethrough(font):
    font._element.set("strike", "sngStrike")


def setHighlight(run, color):
    # get run properties
    rPr = run._r.get_or_add_rPr()

    # Create highlight element
    hl = OxmlElement("a:highlight")

    # Create specify RGB Colour element with color specified
    srgbClr = OxmlElement("a:srgbClr")
    setattr(srgbClr, "val", color)

    # Add colour specification to highlight element
    hl.append(srgbClr)

    # Add highlight element to run properties
    rPr.append(hl)

    return run


def removeBullet(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.insert(
        0,
        etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"),
    )


def removeBullets(textFrame):
    for p in textFrame.paragraphs:
        removeBullet(p)


def removeSelectedBullets(textFrame, removalArray):
    for bulletNumber in removalArray:
        removeBullet(textFrame.paragraphs[bulletNumber])

def findTitleShape(slide):
    if slide.shapes.title == None:
        # Have to use first shape as title
        return slide.shapes[0]
        
    else:
        return slide.shapes.title

def getParagraphs(slide, wantedParagraphs = []):
    paragraphTree = []
    for theShape in slide.shapes:
        if (theShape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX) | (theShape.placeholder_format.type == PP_PLACEHOLDER.OBJECT):
            paragraphTree.append(theShape.text_frame.paragraphs)
    return[ paragraphTree[0][i] for i in wantedParagraphs]
    
def parseText(text):
    textArray = []
    state = "N"
    fragment = ""
    lastChar = ""
    spanState = "None"

    # A span reads a class name or style, then ">", then its text, and an abbr
    # reads a title the same way; both split the accumulated fragment when they
    # close. Emitting a formula part way through would shorten that split and
    # raise IndexError, so formulae stay as source text while one is being read.
    # spanState cannot serve here: the "<span style=" case only sets it when the
    # fragment is already non-empty, so a span at the start of a line misses it.
    structuralFragment = False

    # Lift any inline formulae out of the line before the replacements below can
    # touch them, leaving one sentinel character behind for each. "\[" and "\_"
    # would be rewritten, and the tokeniser drops "*" and "[" outright, so LaTeX
    # left in the line does not reach the fragment it belongs to intact.
    # Testing the line first keeps the cost off the lines that have no formula,
    # and means a presentation without a stylesheet is never told about one.
    mathSpans = []
    mathIndex = 0
    if inlineMathRegex.search(text) and getMathInserter() is not None:
        def stashMath(match):
            if substitutedSentinelRegex.search(match.group(1)):
                return match.group(0)

            mathSpans.append(match.group(1))
            return mathSentinel

        text = inlineMathRegex.sub(stashMath, text)

    # Replace any "\#" strings with entity reference
    text2 = text.replace("\\#", "&#x23;")

    # Replace any "<br/>" strings with newline single character
    text2 = text2.replace("<br/>", "\n")

    # Replace any escaped asterisk strings with entity reference
    text2 = text2.replace("\\*", "&lowast;")

    # Replace any asterisks with spaces either side with entity reference
    text2 = text2.replace(" * ", " &lowast; ")
    if text2[-2:] == " *":
        text2 = text2[:-2] + " &lowast;"

    # Replace any footnote reference starts with char "\uFDD0"
    text2 = text2.replace("[^", u"\uFDD0")

    # Replace any span style starts with char "\uFDD1"
    text2 = re.sub(globals.spanStyleRegex, u"\uFDD1", text2)

    # Replace any span class starts with char "\uFDD2"
    text2 = re.sub(globals.spanClassRegex, u"\uFDD2", text2)

    # Replace any span ends with char "\uFDD3"
    text2 = text2.replace("</span>", u"\uFDD3")

    # Replace any abbreviation starts with char "\uFDD4"
    text2 = text2.replace("<abbr title=", u"\uFDD4")

    # Replace any abbreviation ends with char "\uFDD5"
    text2 = text2.replace("</abbr>", u"\uFDD5")

    # Replace any \[ with char "\uFDD6"
    text2 = text2.replace(r"\[", u"\uFDD6")

    # Replace any \] with char "\uFDD7"
    text2 = text2.replace(r"\]", u"\uFDD7")

    # "\uFDD8" is link separator special character. See below

    # Replace any {~~ with char "\uFDD8"
    text2 = text2.replace("{~~", u"\uFDD8")

    # Replace any ~~} with char "\uFDD8"
    text2 = text2.replace("~~}", u"\uFDD8")

    # Replace any {== with char "\uFDD9"
    # Danish character)
    text2 = text2.replace("{==", u"\uFDD9")

    # Replace any ==} with char "\uFDD9"
    text2 = text2.replace("==}", u"\uFDD9")

    # Replace any {>> with char "\uFDDA"
    text2 = text2.replace("{>>", u"\uFDDA")

    # Replace any <<} with char "\uFDDA"
    text2 = text2.replace("<<}", u"\uFDDA")

    # Replace any {-- with char "\uFDDB"
    text2 = text2.replace("{--", u"\uFDDB")

    # Replace any --} with char "\uFDDB
    text2 = text2.replace("--}", u"\uFDDB")

    # Replace any {++ with char "\uFDDC"
    text2 = text2.replace("{++", u"\uFDDC")

    # Replace any ++} with char "\uFDDC"
    text2 = text2.replace("++}", u"\uFDDC")

    # Replace any <ins> with char "\uFDDD"
    text2 = text2.replace("<ins>", u"\uFDDD")

    # Replace any </ins> with char "\uFDDD"
    text2 = text2.replace("</ins>", u"\uFDDD")

    # Replace any <del> with char "\uFDDE"
    text2 = text2.replace("<del>", u"\uFDDE")

    # Replace any </del> with char "\uFDDE"
    text2 = text2.replace("</del>", u"\uFDDE")

    # Replace any <sub> with char "\uFDDF"
    text2 = text2.replace("<sub>", u"\uFDDF")

    # Replace any </sub> with char "\uFDDF"
    text2 = text2.replace("</sub>", u"\uFDDF")

    # Replace any <sup> with char "\uFDE0"
    text2 = text2.replace("<sup>", u"\uFDE0")

    # Replace any </sup> with char "\uFDE0"
    text2 = text2.replace("</sup>", u"\uFDE0")

    # Note FDE1 - FDE3 used in resolveSymbols
    
    # Handle escaped underscore
    text2 = text2.replace("\\_", "_")

    # Unescape any numeric character references
    text3 = resolveSymbols(text2)

    for c in text3:
        if c == "*":
            # Changing state
            if state == "N":
                # First * potentially starts italic
                textArray.append([state, fragment])
                fragment = ""
                state = "I"

            elif state == "I":
                # Either go to bold or end italic
                if lastChar == "*":
                    # Go to bold
                    state = "B1"

                else:
                    # End italic
                    textArray.append([state, fragment])
                    fragment = ""
                    state = "N"

            elif state == "B1":
                # Starting to close bold bracket
                state = "B2"

            elif lastChar == "*":
                # closing either bold or italic bracket
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif c == "`":
            if state == "N":
                # Going to code
                textArray.append([state, fragment])
                fragment = ""
                state = "C"

            else:
                # exiting code
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif c == u"\uFDD8":
            # Entering or leaving CriticMarkup replacement
            if state == "N":
                # Going to CriticMarkup replacement
                textArray.append([state, fragment])
                fragment = ""
                state = "CMRep"

            else:
                # exiting CriticMarkup replacement
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif c == u"\uFDD9":
            # Entering or leaving CriticMarkup highlight
            if state == "N":
                # Going to CriticMarkup highlight
                textArray.append([state, fragment])
                fragment = ""
                state = "CMHig"

            else:
                # exiting CriticMarkup highlight
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif c == u"\uFDDA":
            # Entering or leaving CriticMarkup comment
            if state == "N":
                # Going to CriticMarkup comment
                textArray.append([state, fragment])
                fragment = ""
                state = "CMCom"

            else:
                # exiting CriticMarkup comment
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif c == u"\uFDDB":
            # Entering or leaving CriticMarkup deletion
            if state == "N":
                # Going to CriticMarkup deletion
                textArray.append([state, fragment])
                fragment = ""
                state = "CMDel"

            else:
                # exiting CriticMarkup deletion
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif c == u"\uFDDC":
            # Entering or leaving CriticMarkup addition
            if state == "N":
                # Going to CriticMarkup addition
                textArray.append([state, fragment])
                fragment = ""
                state = "CMAdd"

            else:
                # exiting CriticMarkup addition
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif c == u"\uFDDD":
            # Entering or leaving underline
            if state == "N":
                # Going to underline
                textArray.append([state, fragment])
                fragment = ""
                state = "Ins"

            else:
                # exiting underline
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif c == u"\uFDDE":
            # Entering or leaving strikethrough
            if state == "N":
                # Going to strikethrough
                textArray.append([state, fragment])
                fragment = ""
                state = "Del"

            else:
                # exiting strikethrough
                textArray.append([state, fragment])
                fragment = ""
                state = "N"
        elif c == u"\uFDDF":
            # Entering or leaving subscript
            if state == "N":
                # Going to subscript
                textArray.append([state, fragment])
                fragment = ""
                state = "Sub"

            else:
                # exiting subscript
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif c == u"\uFDE0":
            # Entering or leaving superscript
            if state == "N":
                # Going to superscript
                textArray.append([state, fragment])
                fragment = ""
                state = "Sup"

            else:
                # exiting superscript
                textArray.append([state, fragment])
                fragment = ""
                state = "N"

        elif c == "[":
            if state == "N":
                # Could be entering a Link
                if fragment != "":
                    textArray.append([state, fragment])

                # The bracket is kept in in case there is no matching ]
                fragment = "["
                state = "LinkText1"

            elif state == "LinkText2":
                # Could be entering an indirect reference
                indLinkText = fragment[:-1]

                # The bracket is kept in in case there is no matching ]
                fragment = "["
                state = "LinkRef1"

        elif c == "]":
            # Could be ending picking up the link text
            if state == "LinkText1":
                # Picked up end of link text
                state = "LinkText2"

                # Remove [ and add a separator to allow for link URL
                fragment = fragment[1:] + u"\uFDD8"

            elif state == "fnref":
                # This terminates a footnote reference
                textArray.append([state, fragment])
                state = "N"
                fragment = ""

            elif state == "LinkRef1":
                # Picked up link reference
                reference = fragment[1:]

                # Attempt to look up reference
                foundReference = False
                for indref, indURL in globals.indirectAnchors:
                    if indref == reference:
                        foundReference = True
                        break

                if foundReference:
                    # Append fragment with resolved reference
                    textArray.append(["Link", indLinkText + u"\uFDD8" + indURL])
                else:
                    print(f"Reference {reference} not resolved.\n")
                    textArray.append(["N", indLinkText])
                fragment = ""
                state = "N"

            else:
                # This was an ordinary square bracket
                fragment += "]"

        elif c == "(":
            # Could be starting to pick up the link URL
            if state == "LinkText2":
                # Picked up start of link URL
                state = "LinkURL1"
            else:
                fragment = fragment + c

        elif c == ")":
            # Could be ending picking up the link URL
            if state == "LinkURL1":
                # Picked up end of link URL
                textArray.append(["Link", fragment])
                fragment = ""
                state = "N"
            else:
                fragment = fragment + c

        elif c == mathSentinel:
            latex = mathSpans[mathIndex]
            mathIndex += 1

            if structuralFragment or state not in mathFragmentStates:
                # The fragment being built carries structure that a formula would
                # break: link text, a footnote reference, a span's class name or
                # a glossary title. Leave the source as text instead.
                fragment = fragment + inlineMathSource(latex)

            else:
                if fragment != "":
                    # Bold accumulates as "B1" and is only emitted as "B2" when
                    # the closing "**" arrives. Flushing it as "B1" would lose
                    # the bold, as addFormattedText has no case for "B1".
                    textArray.append(["B2" if state == "B1" else state, fragment])
                    fragment = ""

                textArray.append(["Math", latex])

        elif c == u"\uFDE3":
            fragment = fragment + "`"

        elif c == u"\uFDE1":
            fragment = fragment + "<"

        elif c == u"\uFDD6":
            fragment = fragment + "["

        elif c == u"\uFDD7":
            fragment = fragment + "]"

        elif c == u"\uFDD5":
            structuralFragment = False
            dictEntry = fragment.split(">")
            dictAbbrev = dictEntry[1]
            dictFull = dictEntry[0].strip().strip("'").strip('"')
            globals.abbrevDictionary[dictAbbrev] = dictFull
            textArray.append(["Gloss", dictAbbrev, dictAbbrev, dictFull])
            fragment = ""

        elif c == u"\uFDD4":
            structuralFragment = True
            if fragment != "":
                textArray.append([state, fragment])
                fragment = ""
            dictEntry = ""

        elif c == u"\uFDD3":
            # End of span
            structuralFragment = False
            if spanState == "Class":
                # Span with class
                splitting = fragment.split(">")
                spanText = splitting[1]
                className = splitting[0].strip().strip("'").strip('"').lower()
                styleText = ""
                if (
                    (className in globals.bgcolors)
                    | (className in globals.fgcolors)
                    | (className in globals.emphases)
                    | (className in globals.fontsizes)
                    | (className in globals.cellcolors)
                    | (className in globals.cellBorderStyling)
                ):
                    textArray.append(["SpanClass", [className, spanText]])

                    fragment = ""

                else:
                    print(
                        f"{className} is not defined. Ignoring reference to it in <span> element."
                    )

                    fragment = spanText
            else:
                # Span with style
                splitting = fragment.split(">")
                spanText = splitting[1]
                styleText = splitting[0].strip().strip("'").strip('"')
                textArray.append(["SpanStyle", [styleText, spanText]])
                className = ""
                fragment = ""

            spanState = "None"

        elif c == u"\uFDD2":
            # In span element where we hit the class name
            structuralFragment = True
            if fragment != "":
                textArray.append([state, fragment])

                fragment = ""
            spanState = "Class"

        elif c == u"\uFDD1":
            # In span element where we hit the style text
            structuralFragment = True
            if fragment != "":
                textArray.append([state, fragment])

                fragment = ""
                spanState = "Style"
        elif c == u"\uFDD0":
            if fragment != "":
                textArray.append([state, fragment])

                fragment = ""
                state = "fnref"

        else:
            fragment = fragment + c

        lastChar = c

    if fragment != "":
        textArray.append([state, fragment])
    return textArray


# Calls the tokeniser and then handles the fragments it gets back
def addFormattedText(p, text):
    boldBold = globals.processingOptions.getCurrentOption("boldBold")
    boldColour = globals.processingOptions.getCurrentOption("boldColour")
    italicItalic = globals.processingOptions.getCurrentOption("italicItalic")
    italicColour = globals.processingOptions.getCurrentOption("italicColour")
    monoFont = globals.processingOptions.getCurrentOption("monoFont")

    # Get back parsed text fragments, along with control information on each
    # fragment
    parsedText = parseText(text)

    # Replace u"\uFDE2" with > in each Fragment
    for f in range(len(parsedText)):
        if parsedText[f][0] in ["SpanClass", "SpanStyle"]:
            parsedText[f][1][1] = parsedText[f][1][1].replace(u"\uFDE2", ">")
        else:
            parsedText[f][-1] = parsedText[f][-1].replace(u"\uFDE2", ">")

    # Prime flattened Text
    flattenedText = ""
    for fragment in parsedText:
        if fragment[0] == "Gloss":
            fragType, fragDetail, fragTerm, fragTitle = fragment
        else:
            fragType, fragDetail = fragment

        # A formula is emitted whole. Splitting it on a newline, or running the
        # entity substitutions below over it, would corrupt its LaTeX.
        if fragType == "Math":
            emitInlineMath(p, fragDetail)
            flattenedText = flattenedText + fragDetail
            continue

        # Break into subfragments around a newline
        if fragType == "SpanClass":
            className, fragText = fragDetail
            styleText = ""
            subfragments = fragText.split("\n")
        elif fragType == "SpanStyle":
            styleText, fragText = fragDetail
            className = ""
            subfragments = fragText.split("\n")
        else:
            subfragments = fragDetail.split("\n")

        # Process each subfragment
        sfnum = 0
        for subfragment in subfragments:
            if sfnum > 0:
                # Subfragments after the first need to be preceded by a line break
                p.add_line_break()

            sfnum += 1
            # Ensure "\*" is rendered as a literal asterisk
            subfragment = subfragment.replace("&lowast;", "*")

            # Ensure "\#" is rendered as a literal octothorpe
            subfragment = subfragment.replace("&#x23;", "#")

            run = p.add_run()

            if fragType not in ["Link", "fnref", "Gloss"]:
                run.text = subfragment
            elif fragType == "Gloss":
                run.text = fragTerm

            if fragType == "I":
                font = run.font

                if italicItalic == True:
                    font.italic = True

                if italicColour != ("None", ""):
                    setColour(font.color, italicColour)

            elif fragType == "Gloss":
                # Add this run to abbrevRunsDictionary - for Glossary fix ups later
                if fragTerm not in globals.abbrevRunsDictionary:
                    globals.abbrevRunsDictionary[fragTerm] = []

                globals.abbrevRunsDictionary[fragTerm].append(run)

            elif fragType == "fnref":
                font = run.font
                font.size = Pt(16)
                setSuperscript(font)
                fnref = fragment[1]
                if fnref in globals.footnoteReferences:
                    footnoteNumber = globals.footnoteReferences.index(fnref)
                    if (
                        globals.processingOptions.getCurrentOption("footnotesOnSlide")
                        == "yes"
                    ):
                        # The definitions are laid out as "[n] text" at the foot
                        # of the slide, so bracket the reference to match
                        run.text = "[" + str(footnoteNumber + 1) + "]"
                    else:
                        run.text = str(footnoteNumber + 1)
                    # Support multiple runs per footnote reference
                    if footnoteNumber not in globals.footnoteRunsDictionary:
                        globals.footnoteRunsDictionary[footnoteNumber] = []
                    globals.footnoteRunsDictionary[footnoteNumber].append(run)

                else:
                    run.text = "[?]"
                    print("Error: Footnote reference '" + fnref + "' unresolved.")

                linkText = "!"
                fragment = ""

            elif fragType == "SpanClass":
                handleSpanClass(run, className)

            elif fragType == "SpanStyle":
                handleSpanStyle(run, styleText)

            elif fragType == "B2":
                font = run.font

                if boldBold == True:
                    font.bold = True

                if boldColour != ("None", ""):
                    setColour(font.color, boldColour)

            elif fragType == "C":
                font = run.font
                font.name = monoFont
            elif fragType == "CMRep":
                font = run.font
                font.color.rgb = RGBColor(255, 140, 0)
                run.text = "{~~" + subfragment + "~~}"
            elif fragType == "CMHig":
                font = run.font
                font.color.rgb = RGBColor(195, 0, 195)
                run.text = "{==" + subfragment + "==}"
            elif fragType == "CMCom":
                font = run.font
                font.color.rgb = RGBColor(0, 0, 195)
                run.text = "{>>" + subfragment + "<<}"
            elif fragType == "CMDel":
                font = run.font
                font.color.rgb = RGBColor(195, 0, 0)
                run.text = "{--" + subfragment + "--}"
            elif fragType == "CMAdd":
                font = run.font
                font.color.rgb = RGBColor(0, 195, 0)
                run.text = "{++" + subfragment + "++}"
            elif fragType == "Ins":
                font = run.font
                font.underline = True
            elif fragType == "Del":
                font = run.font
                setStrikethrough(font)
            elif fragType == "Sub":
                font = run.font
                setSubscript(font)
            elif fragType == "Sup":
                font = run.font
                setSuperscript(font)
            elif fragType == "Link":
                linkArray = subfragment.split(u"\uFDD8")
                linkText = linkArray[0]
                if len(linkArray) > 1:
                    linkURL = linkArray[1]
                else:
                    linkURL = ""
                run.text = linkText
                if linkURL.startswith("#"):
                    # Is an internal Url
                    linkHref = linkURL[1:].strip()
                    # Support multiple runs per href target
                    if linkHref not in globals.href_runs:
                        globals.href_runs[linkHref] = []
                    globals.href_runs[linkHref].append(run)
                else:
                    # Not an internal link so create it
                    hlink = run.hyperlink
                    hlink.address = linkURL

                    # URL might be a macro reference
                    if linkURL[:11] == "ppaction://":
                        # URL is indeed a macro reference, so treat it as such
                        hlink._hlinkClick.action = linkURL

            # Add the flattened text from this subfragment
            if fragType == "Link":
                flattenedText = flattenedText + linkText
            else:
                flattenedText = flattenedText + subfragment

    return flattenedText

def handleSpanClass(run, className):
    if className in globals.bgcolors:
        run = setHighlight(run, globals.bgcolors[className])

    if className in globals.fgcolors:
        font = run.font
        font.color.rgb = RGBColor.from_string(globals.fgcolors[className])

    if className in globals.emphases:
        font = run.font
        if " bold " in " " + globals.emphases[className] + " ":
            font.bold = True
        else:
            font.bold = False
        if " italic " in " " + globals.emphases[className] + " ":
            font.italic = True
        else:
            font.italic = False
        if " underline " in " " + globals.emphases[className] + " ":
            font.underline = True
        else:
            font.underline = False

    if className in globals.fontsizes:
        font = run.font
        font.size = Pt(float(globals.fontsizes[className]))


def handleSpanStyle(run, styleText):
    styleElements = styleText.split(";")

    # Handle the non-empty ones - as the empty one is after the final semicolon
    for styleElement in list(filter(lambda e: e != "", styleElements)):
        styleElementSplit = styleElement.split(":")
        styleElementName = styleElementSplit[0].strip()
        styleElementValue = styleElementSplit[1].strip()

        if styleElementName == "color":
            check, RGBstring = parseRGB(styleElementValue)
            if check:
                run.font.color.rgb = RGBColor.from_string(RGBstring)
            else:
                print(f"Invalid {styleElementName} RGB value {styleElementValue}")

        elif styleElementName == "background-color":
            check, RGBstring = parseRGB(styleElementValue)
            if check:
                setHighlight(run, RGBstring)
            else:
                print(f"Invalid {styleElementName} RGB value {styleElementValue}")

        elif styleElementName == "text-decoration":
            if styleElementValue == "underline":
                run.font.underline = True

        elif styleElementName == "font-weight":
            if styleElementValue == "bold":
                run.font.bold = True

        elif styleElementName == "font-style":
            if styleElementValue == "italic":
                run.font.italic = True

        elif styleElementName == "font-size":
            run.font.size = Pt(float(styleElementValue[:-2]))

